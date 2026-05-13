"""Generate a PowerPoint presentation for the Safe ICU RL project.

This script reads the results CSV, creates a model comparison chart, and
assembles slide content based on the repository structure and README.

Requirements: python-pptx, pandas, matplotlib
Run: python scripts/generate_presentation.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import pandas as pd
import matplotlib.pyplot as plt


ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULTS_DIR = os.path.join(ROOT, "results", "plots")
os.makedirs(RESULTS_DIR, exist_ok=True)


def read_metrics(csv_path):
    if os.path.exists(csv_path):
        return pd.read_csv(csv_path)
    return None


def make_model_comparison_plot(df, outpath):
    plt.figure(figsize=(8, 4))
    x = df['Model']
    y = df['cumulative_reward']
    bars = plt.bar(x, y, color=['#d9534f', '#f0ad4e', '#5cb85c'])
    plt.ylabel('Cumulative Reward')
    plt.title('Model Comparison: Cumulative Reward')
    for bar in bars:
        h = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2, h + 5, f'{h:.1f}', ha='center')
    plt.tight_layout()
    plt.savefig(outpath)
    plt.close()


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    if subtitle:
        subtitle_box.text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0


def add_image_slide(prs, title, image_path, left=Inches(1), top=Inches(1.6), width=Inches(8)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, left, top, width=width)


def add_two_column_slide(prs, title, left_title, left_lines, right_title, right_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    # Use simple layout: join left and right with spacing
    left_block = left_title + '\n' + '\n'.join(left_lines)
    right_block = right_title + '\n' + '\n'.join(right_lines)
    p = tf.paragraphs[0]
    p.text = left_block + '\n\n' + right_block
    p.font.size = Pt(12)


def build_presentation(metrics_df=None):
    prs = Presentation()

    # Title
    add_title_slide(prs, 'Safe Offline POMDP-Based RL for ICU Treatment', 'Final End-to-End Workflow & Results')

    # Agenda
    agenda = [
        'Problem Statement & Objectives',
        'Literature & Gaps',
        'System Architecture & Models',
        'Implementation Details',
        'Results & Analysis',
        'Demo & Dashboard',
        'Conclusions & Future Work'
    ]
    add_bullets_slide(prs, 'Agenda', agenda)

    # Problem Statement
    prob = [
        'Optimize sequential ICU treatment decisions using offline historical data (MIMIC-III).',
        'Minimize patient deterioration while ensuring clinical safety and explainability.',
        'Challenge: Partial observability, distributional risk, and safety constraints.'
    ]
    add_bullets_slide(prs, 'Problem Statement', prob)

    # Final End-to-End Workflow (as bullets per phase)
    workflow = [
        'Phase 1: Data Ingestion & Preprocessing — hourly windowing, imputation, SOFA, reward engineering.',
        'Phase 2: Sequential Representation — LSTM encoder (24h) + VAE → 16D belief state.',
        'Phase 3: Safe Offline RL — Conservative Q-Learning (CQL) with quantile regression (QR-CQL).',
        'Phase 4: Clinical Safety Layer — vitals watchdog + drug interaction vetoes.',
        'Phase 5: Real-time Inference & Dashboard — ensemble uncertainty, SHAP explainability, audit logs.'
    ]
    add_bullets_slide(prs, 'End-to-End Workflow', workflow)

    # System Architecture & Model
    arch_left = ['POMDP formulation', 'Observation: vitals & labs (hourly)', 'Belief: 16D latent vector']
    arch_right = ['Policy: CQL (discrete actions: none/low/med/high)', 'Uncertainty: ensemble heads + QR', 'Safety: rule-based veto']
    add_two_column_slide(prs, 'System Architecture', 'High-level', arch_left, 'Model Components', arch_right)

    # Deep Neural Network Architecture
    dnn = [
        'Temporal Encoder: LSTM (input_dim=16 → hidden=64 → embed)',
        'Latent Model: VAE (latent_dim=16) for belief state',
        'CQL: Conservative Q-learning via d3rlpy (discrete setting)'
    ]
    add_bullets_slide(prs, 'Deep Neural Network Architecture', dnn)

    # Implementation Architecture / Key classes
    impl = [
        'Preprocessing: src.preprocessing.process_data -> MIMICPreprocessor',
        'Temporal: src.temporal_model.lstm_encoder -> TemporalEncoder',
        'Latent: src.latent_model.vae -> LatentVAE',
        'Agents: src.rl_agents.cql_agent -> CQLAgent; baselines in baselines.py',
        'Safety: src.safety.safety_constraints -> SafetyLayer',
        'Explainability: src.explainability.shap_explainer -> ICUExplainer',
        'Dashboard: dashboard/app.py (Streamlit)'
    ]
    add_bullets_slide(prs, 'Implementation - Key Modules', impl)

    # Implementation details / training
    training = [
        'Training orchestration: src.training.train_all.run_pipeline()',
        'Temporal encoder and VAE initialized (hidden_dim=64, latent_dim=16); training skipped in scaffold demo.',
        'CQL configured with batch_size=256 (d3rlpy) — offline fitting on processed parquet trajectories.'
    ]
    add_bullets_slide(prs, 'Implementation - Training', training)

    # Results & Analysis: include chart if metrics available
    if metrics_df is not None:
        chart_path = os.path.join(RESULTS_DIR, 'model_comparison.png')
        make_model_comparison_plot(metrics_df, chart_path)
        add_image_slide(prs, 'Results - Model Comparison', chart_path)

        # Metrics table slide (summary)
        rows = []
        for _, r in metrics_df.iterrows():
            rows.append(f"{r['Model']}: CumR={r['cumulative_reward']:.1f}, AvgR={r['average_reward']:.2f}, SofaΔ={r['sofa_improvement']}")
        add_bullets_slide(prs, 'Results - Key Metrics', rows)
    else:
        add_bullets_slide(prs, 'Results - Key Metrics', ['Result CSV not found.'])

    # Hyperparameters & Tuning (placeholders from repo)
    hp = [
        'TemporalEncoder: hidden_dim=64, num_layers=2',
        'LatentVAE: hidden_dim=64, latent_dim=16',
        'CQL: batch_size=256, n_steps ~ 1k-10k (research)',
        'Optimization: Adam lr=1e-3 for encoder/VAE'
    ]
    add_bullets_slide(prs, 'Training & Hyperparameters', hp)

    # Safety & Explainability
    safety = [
        'SafetyLayer enforces clinical thresholds (BP, HR, SpO2) and drug interaction vetoes.',
        'Every override increments audit counters and is logged in Safety Audit tab.',
        'Explainability via SHAP: per-decision feature impacts saved to results/plots/explainability/.'
    ]
    add_bullets_slide(prs, 'Safety Layer & Explainability', safety)

    # Demo & How to run
    demo = [
        'Run preprocessing: python -m src.preprocessing.process_data',
        'Train models: python -m src.training.train_all',
        'Launch dashboard: streamlit run dashboard/app.py',
        'SHAP explainers and audit logs are generated during inference flows.'
    ]
    add_bullets_slide(prs, 'Demo - How to Run', demo)

    # Discussion, Conclusion, Future Work
    discuss = [
        'Discussion: CQL shows strong cumulative reward in offline benchmarks; safety layer ensures clinical constraints.',
        'Limitations: scaffold demo skips heavy training; real deployment needs prospective validation and clinician-in-the-loop testing.',
        'Future: online fine-tuning (with safeguards), richer drug ontology, patient subgroup analysis.'
    ]
    add_bullets_slide(prs, 'Discussion & Future Work', discuss)

    # References
    refs = [
        'MIMIC-III dataset',
        'Kumar et al., Conservative Q-Learning (CQL)',
        'SHAP: Lundberg & Lee',
        'FQE and OPE literature'
    ]
    add_bullets_slide(prs, 'References', refs)

    out_path = os.path.join(ROOT, 'presentation.pptx')
    prs.save(out_path)
    print(f'Presentation saved to {out_path}')


if __name__ == '__main__':
    metrics_csv = os.path.join(ROOT, 'results', 'metrics', 'model_comparison.csv')
    df = read_metrics(metrics_csv)
    build_presentation(df)
